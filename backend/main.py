from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import JSONResponse
from typing import Optional
from schemas import (
    CourseStructureRequest,
    GenerateTitleRequest, GenerateTitleResponse, FetchWebRequest, FetchYouTubeRequest,
    GenerateOutlineBaseRequest, ExportChapterRequest, GenerateVoiceScriptReq, GenerateFlashcardsRequest,
    GenerateMCQRequest, GenerateAssessmentRequest, ChatRequest, ThemeUploadRequest,
    ChatbotBuilderRequest
)
from course_planner import generate_course_structure
from content_generator import generate_chapter_content, generate_course_quiz

from rag_pipeline import ingest_document
from course_store import save_course, get_courses, get_course, update_course, delete_course
from database import save_course_to_mysql
import os
import shutil
import json
from dotenv import load_dotenv
import logging
from openai import OpenAI

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("main")

load_dotenv()

openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
LLM_MODEL = os.getenv("LLM_MODEL", "gpt-4o-mini")

from chatbot_builder_service import parse_number_from_text

app = FastAPI()

from fastapi.exceptions import RequestValidationError
@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request, exc):
    logger.error(f"[Validation Error] Request path: {request.url.path}, Error: {exc.errors()}, Body: {await request.body()}")
    return JSONResponse(
        status_code=422,
        content={"detail": exc.errors(), "body": str(await request.body())}
    )

# Register online course generator router
from app.api.online_course_generator import router as ocg_router
app.include_router(ocg_router)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create uploads directory if not exists
os.makedirs("uploads", exist_ok=True)
app.mount("/uploads", StaticFiles(directory="uploads"), name="uploads")

def validate_uploaded_file(file: UploadFile, max_size_mb: float, allowed_extensions: list[str]):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in allowed_extensions:
        raise HTTPException(
            status_code=400, 
            detail=f"File extension '{ext}' is not allowed. Supported formats: {', '.join(allowed_extensions)}"
        )
    
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    max_size_bytes = max_size_mb * 1024 * 1024
    if file_size > max_size_bytes:
        raise HTTPException(
            status_code=400,
            detail=f"File size exceeds the limit of {max_size_mb}MB. Got {file_size / (1024 * 1024):.2f}MB."
        )

@app.get("/")
def read_root():
    return {"message": "AI Course Builder API is running"}

@app.post("/course/upload-doc")
async def upload_doc(file: UploadFile = File(...)):
    validate_uploaded_file(file, 10.0, ['.pdf', '.docx', '.txt', '.pptx', '.doc'])
    temp_file = file.filename
    try:
        with open(temp_file, "wb") as f:
            shutil.copyfileobj(file.file, f)
        
        chunks_count = ingest_document(temp_file, file.filename)
        
        if os.path.exists(temp_file):
            os.remove(temp_file)
            
        return {"status": "success", "chunks_processed": chunks_count}
    except Exception as e:
        if 'temp_file' in locals() and os.path.exists(temp_file):
            os.remove(temp_file)
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/structure")
async def create_structure(req: CourseStructureRequest):
    structure = generate_course_structure(
        req.details.courseName,
        req.details.description,
        req.details.subject,
        req.details.level
    )
    return {"status": "success", "data": structure}

async def auto_generate_dalle_image(chapter_title: str, explanation: str) -> Optional[str]:
    import requests
    import uuid
    import os
    try:
        visual_prompt = (
            f"A clean professional educational illustration, workflow flow chart, or concept diagram explaining '{chapter_title}'. "
            f"Style: clean modern infographic, dark theme, high vector graphics, highly informative, no text typos."
        )
        response = openai_client.images.generate(
            model="gpt-image-2",
            prompt=visual_prompt[:1000],
            n=1,
            size="1024x1024",
            quality="low"
        )
        logger.info(f"Image response: {response}")
        img_data = None
        first_item = response.data[0]
        
        url_val = getattr(first_item, 'url', None)
        b64_val = getattr(first_item, 'b64_json', None)
        
        if b64_val:
            import base64
            img_data = base64.b64decode(b64_val)
        elif url_val and url_val != "None":
            img_data = requests.get(url_val, timeout=20).content
        else:
            raise ValueError(f"No valid image URL or b64_json found in response: {first_item}")

        unique_filename = f"dalle_{uuid.uuid4().hex[:8]}.png"
        os.makedirs("uploads", exist_ok=True)
        file_path = os.path.join("uploads", unique_filename)
        with open(file_path, "wb") as f:
            f.write(img_data)
        
        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        return f"{base_url}/uploads/{unique_filename}"
    except Exception as e:
        logger.error(f"Failed to auto generate DALL-E visual for {chapter_title}: {e}")
        return None



@app.post("/course/create")
async def finalize_course(course: dict):
    from uuid import uuid4
    cid = course.get("id") or str(uuid4())
    try:
        logger.warning("Saving course to deprecated JSON file storage.")
        try:
            save_course(cid, course)
        except Exception as json_err:
            logger.error(f"Error saving to JSON storage: {json_err}")
            
        # Also save to MySQL and capture the generated ID
        try:
            mysql_course_id = save_course_to_mysql(course)
        except Exception as db_err:
            logger.error(f"Course saved to JSON but MySQL sync failed: {db_err}")
            raise HTTPException(status_code=500, detail=f"Database error: {str(db_err)}")
        return {"status": "success", "course_id": cid, "mysql_course_id": mysql_course_id}

    except ValueError as ve:
        raise HTTPException(status_code=400, detail=str(ve))

@app.put("/course/{course_id}")
async def edit_course(course_id: str, course: dict):
    try:
        logger.warning("Updating course in deprecated JSON file storage.")
        try:
            update_course(course_id, course)
        except Exception as json_err:
            logger.error(f"Error updating JSON storage: {json_err}")
            
        # Also update MySQL
        try:
            save_course_to_mysql(course)
        except Exception as db_err:
            logger.error(f"Course updated in JSON but MySQL sync failed: {db_err}")
            
        return {"status": "success", "course_id": course_id}
    except ValueError as ve:
        raise HTTPException(status_code=400, detail=str(ve))

@app.delete("/course/{course_id}")
async def remove_course(course_id: str):
    delete_course(course_id)
    return {"status": "success"}

@app.get("/courses")
async def list_courses():
    try:
        from database import get_courses_from_mysql
        courses = get_courses_from_mysql()
        # Merge in any locally-saved JSON courses that are not in MySQL
        try:
            json_courses = get_courses()
            seen_ids = {str(c.get("id")) for c in courses if c.get("id") is not None}
            for jc in json_courses:
                if str(jc.get("id")) not in seen_ids:
                    courses.append(jc)
        except Exception:
            pass
        return {"courses": courses}
    except Exception as e:
        logger.error(f"Error fetching from MySQL: {e}")
        logger.warning("MySQL failed. Falling back to deprecated courses.json storage.")
        courses = get_courses()
        return {"courses": courses}

@app.get("/subjects")
async def list_subjects():
    try:
        from database import get_all_subjects_from_mysql
        subjects = get_all_subjects_from_mysql()
        return {"status": "success", "subjects": subjects}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/course/themes")
def get_themes():
    themes_file = "themes.json"
    if not os.path.exists(themes_file):
        return []
    try:
        with open(themes_file, "r", encoding="utf-8") as f:
            import json
            return json.load(f)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading themes: {str(e)}")

@app.post("/course/theme")
def save_theme(req: ThemeUploadRequest):
    themes_file = "themes.json"
    
    allowed_keys = {
        "--bg-primary", "--bg-secondary", "--text-main", "--text-secondary", 
        "--text-muted", "--border-color", "--accent-color", "--accent-bg", 
        "--code-bg", "--code-text", "--theme-shadow",
        "--font-family", "--font-size-base", "--font-size-h1", "--font-size-h2",
        "--font-size-h3", "--line-height", "--block-spacing"
    }
    
    if not req.variables:
        raise HTTPException(status_code=400, detail="Variables mapping cannot be empty.")
        
    for k, v in req.variables.items():
        if k not in allowed_keys:
            raise HTTPException(status_code=400, detail=f"Invalid variable key '{k}'. Allowed keys: {', '.join(allowed_keys)}")
            
        v_str = str(v).strip()
        if ";" in v_str or "}" in v_str or "<" in v_str or ">" in v_str:
             raise HTTPException(status_code=400, detail=f"Invalid characters in value for key '{k}'.")
    
    themes_list = []
    if os.path.exists(themes_file):
        try:
            with open(themes_file, "r", encoding="utf-8") as f:
                import json
                themes_list = json.load(f)
        except Exception:
            themes_list = []
            
    themes_list = [t for t in themes_list if t.get("id") != req.id]
    
    new_theme = {
        "id": req.id,
        "name": req.name,
        "variables": req.variables
    }
    themes_list.append(new_theme)
    
    try:
        with open(themes_file, "w", encoding="utf-8") as f:
            import json
            json.dump(themes_list, f, indent=2)
        return {"status": "success", "theme_id": req.id}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving theme: {str(e)}")

@app.get("/course/{course_id}")
async def get_single_course(course_id: str):
    try:
        from database import get_course_details_from_mysql
        # Convert string ID to int if it's numeric
        cid = int(course_id) if course_id.isdigit() else None
        
        if cid:
            course = get_course_details_from_mysql(cid)
            if course:
                return {"status": "success", "course": course}
        
        # Fallback to JSON if not found in MySQL
        logger.warning("Course not found in MySQL. Falling back to deprecated JSON file storage lookup.")
        from course_store import get_course
        course = get_course(course_id)
        if course:
            return {"status": "success", "course": course}
            
        raise HTTPException(status_code=404, detail="Course not found")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/generate-title", response_model=GenerateTitleResponse)
async def api_generate_title(req: GenerateTitleRequest):
    from content_generator import generate_brief_title
    title = generate_brief_title(req.description)
    return {"title": title}

@app.post("/course/upload-thumbnail")
async def upload_thumbnail(file: UploadFile = File(...)):
    validate_uploaded_file(file, 5.0, ['.jpg', '.jpeg', '.png', '.gif', '.webp'])
    import uuid
    import shutil
    
    extension = os.path.splitext(file.filename)[1]
    unique_filename = f"{uuid.uuid4()}{extension}"
    file_path = os.path.join("uploads", unique_filename)
    
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        return {"status": "success", "url": f"{base_url}/uploads/{unique_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/fetch-web")
async def fetch_web(req: FetchWebRequest):
    import requests
    from bs4 import BeautifulSoup
    import uuid
    
    try:
        response = requests.get(req.url, timeout=10)
        soup = BeautifulSoup(response.content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        
        unique_filename = f"{uuid.uuid4()}_web.txt"
        file_path = os.path.join("uploads", unique_filename)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)
            
        chunks_count = ingest_document(file_path, unique_filename)
        return {"status": "success", "chunks_processed": chunks_count}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/fetch-youtube")
async def fetch_youtube(req: FetchYouTubeRequest):
    from youtube_transcript_api import YouTubeTranscriptApi
    import urllib.parse
    import uuid
    
    try:
        parsed_url = urllib.parse.urlparse(req.youtube_url)
        video_id = urllib.parse.parse_qs(parsed_url.query).get("v")
        if not video_id:
            # Handle youtu.be format
            video_id = [parsed_url.path.lstrip("/")]
        
        if not video_id:
            raise Exception("Invalid YouTube URL")
            
        transcript = YouTubeTranscriptApi.get_transcript(video_id[0])
        text = " ".join([t['text'] for t in transcript])
        
        unique_filename = f"{uuid.uuid4()}_youtube.txt"
        file_path = os.path.join("uploads", unique_filename)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(text)
            
        chunks_count = ingest_document(file_path, unique_filename)
        return {"status": "success", "chunks_processed": chunks_count}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/generate-outline")
async def generate_outline(req: GenerateOutlineBaseRequest):
    from content_generator import generate_outline_skeleton
    structure = generate_outline_skeleton(req.description, req.modules_count, req.chapters_per_module)
    return {"status": "success", "data": structure}

@app.post("/course/export")
async def export_chapter(req: ExportChapterRequest):
    import uuid
    unique_filename = f"export_{uuid.uuid4()}"
    file_path = os.path.join("uploads", unique_filename)
    
    try:
        if req.format == "txt":
            file_path += ".txt"
            with open(file_path, "w", encoding="utf-8") as f:
                f.write(req.content.get("explanation", ""))
        elif req.format == "pdf":
            import pdfkit
            file_path += ".pdf"
            html_content = f"<h1>{req.chapter_title}</h1><p>{req.content.get('explanation', '')}</p>"
            try:
                pdfkit.from_string(html_content, file_path)
            except OSError:
                file_path = file_path.replace(".pdf", ".txt")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(f"{req.chapter_title}\n\n{req.content.get('explanation', '')}")
        elif req.format == "pptx":
            from pptx import Presentation
            file_path += ".pptx"
            prs = Presentation()
            slide_layout = prs.slide_layouts[1] # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = req.chapter_title
            content.text = req.content.get("explanation", "")[:500] + "..." # basic
            prs.save(file_path)
        else:
            raise Exception("Unsupported format")

        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        return {"status": "success", "url": f"{base_url}/uploads/{unique_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/voice")
async def generate_voice(req: GenerateVoiceScriptReq):
    import uuid
    
    unique_filename = f"tts_{uuid.uuid4()}.mp3"
    file_path = os.path.join("uploads", unique_filename)
    
    try:
        response = openai_client.audio.speech.create(
            model="tts-1",
            voice="alloy",
            input=req.text
        )
        response.stream_to_file(file_path)
        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        return {"status": "success", "audio_url": f"{base_url}/uploads/{unique_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/flashcards")
async def generate_flashcards(req: GenerateFlashcardsRequest):
    from pydantic import BaseModel
    class FlashcardModel(BaseModel):
        question: str
        answer: str
    class FlashcardsResponse(BaseModel):
        flashcards: list[FlashcardModel]
    
    try:
        completion = openai_client.beta.chat.completions.parse(
            model=LLM_MODEL,
            messages=[
                {"role": "system", "content": "You are a helpful educational AI. Generate exactly 5 flashcards from the provided text."},
                {"role": "user", "content": f"Text:\n{req.text[:2000]}"}
            ],
            response_format=FlashcardsResponse
        )
        return {"status": "success", "flashcards": completion.choices[0].message.parsed.model_dump().get("flashcards")}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def convert_to_pdf_if_needed(file_path: str) -> str:
    import platform
    import subprocess
    
    ext = os.path.splitext(file_path)[1].lower()
    if ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']:
        # Find LibreOffice binary
        libreoffice_bin = "soffice"  # Default in Linux/AWS PATH
        if platform.system() == "Windows":
            # Common Windows paths for LibreOffice
            possible_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            for p in possible_paths:
                if os.path.exists(p):
                    libreoffice_bin = p
                    break
        
        # Output directory is 'uploads'
        out_dir = "uploads"
        
        # Execute headless LibreOffice conversion
        cmd = [libreoffice_bin, "--headless", "--convert-to", "pdf", "--outdir", out_dir, file_path]
        try:
            logger.info(f"Converting {file_path} to PDF using command: {' '.join(cmd)}")
            subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
            
            # The converted file has the same name but with .pdf extension
            pdf_filename = os.path.splitext(os.path.basename(file_path))[0] + ".pdf"
            pdf_path = os.path.join(out_dir, pdf_filename)
            
            if os.path.exists(pdf_path):
                # Delete original file
                if os.path.exists(file_path):
                    os.remove(file_path)
                return pdf_path
        except Exception as err:
            logger.error(f"LibreOffice conversion failed: {err}")
            # If conversion fails, keep original file
            
    return file_path

@app.post("/course/upload-media")
async def upload_media(file: UploadFile = File(...)):
    validate_uploaded_file(file, 10.0, ['.jpg', '.jpeg', '.png', '.gif', '.webp', '.mp4', '.mp3', '.wav', '.pdf', '.docx', '.txt', '.pptx', '.xlsx', '.xls', '.ppt', '.doc'])
    import uuid
    import shutil
    
    extension = os.path.splitext(file.filename)[1]
    unique_filename = f"{uuid.uuid4()}{extension}"
    file_path = os.path.join("uploads", unique_filename)
    
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # Convert Office documents to PDF so they can open securely in SecureDocViewer
        final_file_path = convert_to_pdf_if_needed(file_path)
        final_filename = os.path.basename(final_file_path)
        
        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        return {"status": "success", "url": f"{base_url}/uploads/{final_filename}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/course/list-media")
async def list_media():
    try:
        uploads_dir = "uploads"
        if not os.path.exists(uploads_dir):
            return {"status": "success", "files": []}
        
        base_url = os.getenv("PUBLIC_ASSET_URL", "http://localhost:8000")
        files_list = []
        for filename in os.listdir(uploads_dir):
            file_path = os.path.join(uploads_dir, filename)
            if os.path.isfile(file_path):
                stat = os.stat(file_path)
                files_list.append({
                    "filename": filename,
                    "url": f"{base_url}/uploads/{filename}",
                    "size": stat.st_size,
                    "created_at": stat.st_mtime
                })
        
        files_list.sort(key=lambda x: x["created_at"], reverse=True)
        return {"status": "success", "files": files_list}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/mcq")
async def generate_mcqs(req: GenerateMCQRequest):
    try:
        from schemas import MCQResponse
        prompt = f"""You are an expert educator. Generate exactly 5 multiple-choice questions (each with 4 options) that test understanding for the following context.
Course: {req.course_title}
Module: {req.module_title}
Chapter: {req.chapter_title or 'General'}
Assessment Guidelines: {req.assessment_text or 'None'}"""

        completion = openai_client.beta.chat.completions.parse(
            model=LLM_MODEL,
            messages=[
                {"role": "system", "content": prompt},
            ],
            response_format=MCQResponse
        )
        return {"status": "success", "mcqs": completion.choices[0].message.parsed.model_dump().get("mcqs")}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/assessment")
async def generate_assessment(req: GenerateAssessmentRequest):
    try:
        from schemas import MCQResponse
        prompt = f"""You are an expert educator. Generate exactly 10 multiple-choice questions (each with 4 options) that form a comprehensive assessment for the following module.
Course: {req.course_title}
Module: {req.module_title}
Assessment Guidelines: {req.assessment_text or 'None'}"""

        completion = openai_client.beta.chat.completions.parse(
            model=LLM_MODEL,
            messages=[
                {"role": "system", "content": prompt},
            ],
            response_format=MCQResponse
        )
        return {"status": "success", "mcqs": completion.choices[0].message.parsed.model_dump().get("mcqs")}
    except Exception as e:
        logger.error(f"Error generating assessment: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/course/auto-fill")
async def api_auto_fill():
    from content_generator import generate_course_details_from_context
    details = generate_course_details_from_context()
    if not details:
        raise HTTPException(status_code=400, detail="No source documents found. Please upload documents in Step 1 first.")
    return {"status": "success", "details": details}

@app.post("/course/chat")
async def api_chat(req: ChatRequest):
    from chat_service import build_system_prompt, parse_metadata

    logger.info(f"Chat request received. Scope: {req.scope}, Details: {req.details}")

    system_prompt = build_system_prompt(
        scope=req.scope,
        details=req.details,
        course_data=req.courseData,
        available_subjects=req.availableSubjects
    )

    messages = [{"role": "system", "content": system_prompt}] + req.messages

    try:
        response = openai_client.chat.completions.create(
            model=LLM_MODEL,
            messages=messages,
            temperature=0.7,
            max_tokens=6000
        )
        
        # Log token usage
        if hasattr(response, "usage") and response.usage:
            logger.info(f"Token usage - Prompt: {response.usage.prompt_tokens}, Completion: {response.usage.completion_tokens}, Total: {response.usage.total_tokens}")

        ai_reply = response.choices[0].message.content

        reply_text, metadata, type_val = parse_metadata(
            ai_reply=ai_reply,
            scope=req.scope,
            details=req.details
        )

        return {
            "status": "success",
            "reply": reply_text,
            "metadata": metadata,
            "type": type_val
        }
    except Exception as e:
        logger.error(f"Error in chat completion API: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/course/chatbot-builder/chat")
async def api_chatbot_builder_chat(req: ChatbotBuilderRequest):
    from chatbot_builder_service import build_builder_system_prompt, parse_quick_replies
    from chat_service import parse_metadata

    logger.info(f"Chatbot Builder request received. Step: {req.currentStep}")

    # -------------------------------------------------------------------------
    # OUTLINE_EDIT: Use a dedicated JSON-mode call to guarantee structure output
    # -------------------------------------------------------------------------
    is_outline_edit_req = False
    if req.currentStep in ["OUTLINE_EDIT", "EDIT_OUTLINE_CHOICE", "ASK_REDUCE_COUNT", "ASK_ADD_TOPIC"]:
        user_message = ""
        if req.messages:
            for msg in reversed(req.messages):
                if msg.get("role") == "user":
                    user_message = msg.get("content", "")
                    break
        lowercase_msg = user_message.lower()
        confirm_words = ["yes", "continue", "looks good", "proceed", "generate", "correct", "confirm", "happy", "fine", "ok", "go ahead"]
        
        # Check if the user is confirming a pending addition from a previous 10-module limit warning
        is_pending_add_confirm = False
        if req.messages:
            last_assistant_msg = ""
            for msg in reversed(req.messages):
                if msg.get("role") == "assistant":
                    last_assistant_msg = msg.get("content", "").lower()
                    break
            if "maximum limit of 10" in last_assistant_msg or "i can add" in last_assistant_msg:
                is_pending_add_confirm = any(w in lowercase_msg for w in confirm_words)

        is_confirmation = any(w in lowercase_msg for w in confirm_words) and not any(neg in lowercase_msg for neg in ["not", "dont", "change", "add", "remove", "delete", "reduce"])
        if is_pending_add_confirm:
            is_confirmation = False

        # Check if we should change details (go-back to CONFIRM_DETAILS)
        is_details_redirect = any(w in lowercase_msg for w in ["change", "edit", "modify", "update", "correct"]) and any(w in lowercase_msg for w in ["detail", "details", "topic", "goal", "style", "level", "duration", "objective", "requirements", "hours", "basic info", "info", "basic"])

        # Programmatic Ambiguity Checks
        is_ambiguous = False
        ambiguity_reply = ""
        
        if not is_confirmation:
            has_edit_verb = any(w in lowercase_msg for w in ["edit", "change", "modify", "update", "adjust"])
            has_detail_noun = any(w in lowercase_msg for w in ["detail", "details", "topic", "goal", "style", "level", "duration", "objective", "requirements", "hours", "basic info", "info", "basic"])
            has_structure_noun = any(w in lowercase_msg for w in ["module", "modules", "chapter", "chapters", "outline", "syllabus", "roadmap", "lesson", "lessons"])
            has_action_verb = any(w in lowercase_msg for w in ["add", "remove", "delete", "reduce", "shrink", "decrease", "cut", "rename", "reorder", "shuffle", "move", "swap"])
            
            if has_edit_verb and not has_detail_noun and not has_structure_noun:
                is_ambiguous = True
                ambiguity_reply = "Would you like to modify the course details (such as topic, goal, level, style, or duration) or make changes to the course outline modules?"
            elif has_structure_noun and not has_action_verb and not has_edit_verb:
                is_ambiguous = True
                ambiguity_reply = "I'm not sure if you want to add modules, remove modules, or rename them. Could you please clarify your request?"

        if is_ambiguous:
            return JSONResponse({
                "status": "success",
                "reply": ambiguity_reply,
                "quickReplies": ["Edit course details", "Add new module", "Reduce modules"],
                "metadata": {
                    "next_step": "OUTLINE_EDIT"
                },
                "type": "details"
            })

        # No clarifying redirect since we directly edit the outline without asking questions
        if not is_confirmation and not is_details_redirect:
            if lowercase_msg.strip() != "edit outline":
                is_outline_edit_req = True

    if is_outline_edit_req:
        try:
            current_structure = req.courseData.get("structure", {})
            details = req.courseData.get("details", {})
            
            # Find the user's last edit request
            user_message = "Modify the outline"
            if req.messages:
                for msg in reversed(req.messages):
                    if msg.get("role") == "user":
                        user_message = msg.get("content", "")
                        break
            lowercase_msg = user_message.lower()

            # Custom modification prompt builder for clarifying steps
            current_count = len(current_structure.get("modules", []))
            
            # Check if they specified a count directly in the command
            count = parse_number_from_text(lowercase_msg)
            
            # If confirming pending add, override count
            import re
            if is_pending_add_confirm and req.messages:
                last_assistant_msg = ""
                for msg in reversed(req.messages):
                    if msg.get("role") == "assistant":
                        last_assistant_msg = msg.get("content", "").lower()
                        break
                match = re.search(r'add (\d+) module', last_assistant_msg)
                if match:
                    count = int(match.group(1))
            
            # Determine if this is a specific module/position edit
            is_specific_edit = False
            
            # 1. Match patterns like "module 3", "modules 1 and 3", "chapter 2", etc.
            if re.search(r'\b(modules?|chapters?)\s+(\d+|one|two|three|four|five|six|seven|eight|nine|ten)\b', lowercase_msg):
                is_specific_edit = True
                
            # 2. Check for ordinal/positional descriptors like "first", "3rd", "last"
            ordinals = ["first", "second", "third", "fourth", "fifth", "sixth", "seventh", "eighth", "ninth", "tenth", "last",
                        "1st", "2nd", "3rd", "4th", "5th", "6th", "7th", "8th", "9th", "10th"]
            if any(rf"\b{ord}\b" in lowercase_msg for ord in ordinals):
                is_specific_edit = True
                
            # 3. Check for position indicators like "at 4", "at position 3", "at index 2", "after module 2"
            if any(w in lowercase_msg for w in ["at", "position", "index", "place", "after", "before"]):
                if re.search(r'\b(at|position|index|place|after|before)\s+(\d+|one|two|three|four|five|six|seven|eight|nine|ten)\b', lowercase_msg):
                    is_specific_edit = True
                    
            # 4. Check for multiple indices like "1 and 3", "2 and 4"
            if re.search(r'\b\d+\s+and\s+\d+\b', lowercase_msg):
                is_specific_edit = True
            
            # Check reduce vs add
            is_reduce_req = False
            is_add_req = False
            if not is_specific_edit:
                is_reduce_req = any(w in lowercase_msg for w in ["reduce", "shrink", "delete", "remove", "cut", "decrease"])
                is_add_req = ("add" in lowercase_msg and any(w in lowercase_msg for w in ["module", "modules", "syllabus", "roadmap"])) or is_pending_add_confirm
            
            if is_add_req:
                if count is None:
                    count = 1
                
                # Naming/topic extraction helper
                topic_focus = "a relevant topic of your choice"
                for prep in ["focused on", "focusing on", "focus on", "named", "called", "about", "on", "for"]:
                    if f" {prep} " in f" {lowercase_msg} ":
                        idx = lowercase_msg.find(prep)
                        topic_focus = user_message[idx + len(prep):].strip()
                        topic_focus = topic_focus.strip('."\'? ')
                        break
                        
                target_count = current_count + count
                if target_count > 10:
                    allowable_add = 10 - current_count
                    if allowable_add <= 0:
                        return JSONResponse({
                            "status": "success",
                            "reply": "Your course already has 10 modules, which is the maximum limit. I cannot add any more modules.",
                            "metadata": {
                                "next_step": "OUTLINE_EDIT"
                            },
                            "type": "details"
                        })
                    else:
                        return JSONResponse({
                            "status": "success",
                            "reply": f"I can only create up to 10 modules. I can add {allowable_add} modules to bring your course to the maximum limit of 10. Would you like me to do that?",
                            "quickReplies": ["Yes, add modules", "Cancel"],
                            "metadata": {
                                "next_step": "OUTLINE_EDIT"
                            },
                            "type": "details"
                        })
                else:
                    user_message = f"Add exactly {count} new modules focused on '{topic_focus}' to the outline. The final outline MUST contain exactly {target_count} modules."
                    lowercase_msg = user_message.lower()
                    
            elif is_reduce_req:
                if count is None:
                    count = 1
                    
                is_reduce_to = "to " in lowercase_msg or "reduce to" in lowercase_msg
                if is_reduce_to:
                    target_count = max(0, count)
                    diff = max(0, current_count - target_count)
                else:
                    diff = count
                    target_count = max(0, current_count - diff)
                    
                if target_count == current_count:
                    return JSONResponse({
                        "status": "success",
                        "reply": f"The course outline already contains exactly {current_count} modules. I cannot reduce it to {target_count}.",
                        "metadata": {
                            "next_step": "OUTLINE_EDIT"
                        },
                        "type": "details"
                    })
                elif target_count < 1:
                    return JSONResponse({
                        "status": "success",
                        "reply": "A course outline must contain at least 1 module. I cannot reduce the course below 1 module.",
                        "metadata": {
                            "next_step": "OUTLINE_EDIT"
                        },
                        "type": "details"
                    })
                else:
                    user_message = f"Reduce exactly {diff} modules of your choice from the outline. The final outline MUST contain exactly {target_count} modules."
                    lowercase_msg = user_message.lower()

            # 1. Programmatic Shuffle Interceptor (only for random shuffle requests)
            is_random_shuffle = ("shuffle" in lowercase_msg or "reorder" in lowercase_msg) and not is_specific_edit and any(w in lowercase_msg for w in ["random", "randomly", "shuffle", "mix"])
            if is_random_shuffle:
                modules = current_structure.get("modules", [])
                import random
                if "submodule" in lowercase_msg or "chapter" in lowercase_msg:
                    # Shuffle chapters inside modules
                    for m in modules:
                        if "chapters" in m and isinstance(m["chapters"], list):
                            random.shuffle(m["chapters"])
                    return JSONResponse({
                        "status": "success",
                        "reply": "I have randomly shuffled the chapters (submodules) inside all modules for you. How does the new order look?",
                        "metadata": {
                            "next_step": "OUTLINE_EDIT",
                            "modules": modules
                        },
                        "type": "structure"
                    })
                else:
                    # Shuffle main modules list
                    random.shuffle(modules)
                    return JSONResponse({
                        "status": "success",
                        "reply": "I have randomly shuffled the order of the modules for you. How does the new order look?",
                        "metadata": {
                            "next_step": "OUTLINE_EDIT",
                            "modules": modules
                        },
                        "type": "structure"
                    })

            specific_edit_rule = ""
            if is_specific_edit:
                specific_edit_rule = "\n7. SPECIFIC EDIT RULE: The user has requested a specific modification (e.g., removing specific modules/chapters, inserting a module at a specific index, or reordering). You MUST perform that exact modification, adding/removing exactly those elements, while leaving the rest of the syllabus structure and content completely unchanged."

            edit_prompt = f"""You are an expert curriculum designer.
Your task is to modify the current course outline based strictly on the user's request.

Course Details:
- Name: {details.get("courseName")}
- Subject: {details.get("subject")}
- Goal: {details.get("description")}
- Level: {details.get("level")}
- Duration: {details.get("duration")} Hours

Current Course Outline (JSON):
{json.dumps(current_structure, indent=2)}

User's Modification Request:
"{user_message}"

Rules:
1. STRICT COMPLETE OUTPUT RULE: You MUST output the ENTIRE updated course syllabus JSON including all unchanged modules. Never omit, truncate, or drop any module or chapter from the current outline.
2. CHAPTERS GENERATION RULE: Every module in the output (including newly added ones) MUST contain a list of relevant chapters. The 'chapters' list must NEVER be empty. If a new module is added, you MUST generate at least 3-4 relevant chapters (subtopics) for it.
3. EXACT COUNT RULE: The final outline in the JSON output MUST have EXACTLY the target count of modules specified in the user request description. Count them carefully. If it says "The final outline MUST contain exactly X modules", your JSON output 'modules' array MUST have exactly X items. Do not output X-1 or X+1 modules.
4. RENAME STRUCTURAL PRESERVATION RULE: If the request is to rename modules or chapters, you MUST keep the exact same number of modules and chapters, and only change/rename their titles. Do NOT alter the course structure or count of modules/chapters.
5. NO INDEXES RULE: Do NOT prepend numbers, chapter numbers, or index prefixes (like "Module 1", "Chapter 1 -", "1.1") to module or chapter titles.
6. JSON ONLY RULE: Output ONLY valid JSON conforming to the schema. No markdown code blocks, no text explanations.{specific_edit_rule}

Expected JSON output format exactly:
{{
  "next_step": "OUTLINE_EDIT",
  "modules": [
    {{
      "title": "Module Title",
      "chapters": [
        {{"title": "Chapter Title"}},
        {{"title": "Another Chapter Title"}}
      ]
    }}
  ]
}}"""

            response = openai_client.chat.completions.create(
                model=LLM_MODEL,
                messages=[
                    {"role": "system", "content": "You are a curriculum design expert. Output only valid JSON conforming to the schema. No markdown, no extra text."},
                    {"role": "user", "content": edit_prompt}
                ],
                temperature=0.1,
                max_tokens=4000,
                response_format={"type": "json_object"}
            )

            from metering_helper import track_chatbot_cost
            track_chatbot_cost(req.draft_id, response, LLM_MODEL, "outline_edit")

            structure_json = json.loads(response.choices[0].message.content)
            structure_json["next_step"] = "OUTLINE_EDIT"

            return {
                "status": "success",
                "reply": "Here is the updated course structure outline. Please take a moment to review it. Would you like to make any further modifications, or are you happy with this outline?",
                "quickReplies": ["Confirm Outline", "Reduce modules", "Add new module", "Rename modules/chapters"],
                "metadata": structure_json,
                "type": "structure"
            }
        except Exception as e:
            logger.error(f"Error in OUTLINE_EDIT JSON modification: {str(e)}")
            raise HTTPException(status_code=500, detail=str(e))

    # -------------------------------------------------------------------------
    # All other steps: standard conversational slot-filling flow
    # -------------------------------------------------------------------------
    from chatbot_builder_service import (
        extract_slots_from_message,
        determine_next_step,
        build_builder_system_prompt,
        parse_quick_replies,
        reinject_quick_replies_into_history
    )
    from chat_service import parse_metadata

    try:
        # Extract user last message
        user_message = ""
        if req.messages:
            for msg in reversed(req.messages):
                if msg.get("role") == "user":
                    user_message = msg.get("content", "")
                    break

        # Map legacy details key/values to clean slot values
        details = req.courseData.get("details", {}) or {}
        current_slots = {
            "topic": details.get("topic") or details.get("subject") or details.get("courseName") or "",
            "learningGoal": details.get("learningGoal") or details.get("description") or "",
            "currentLevel": details.get("currentLevel") or details.get("level") or "",
            "learningStyle": details.get("learningStyle") or details.get("requirements") or "",
            "duration": details.get("duration") or "",
            "language": "English"
        }

        # If the user is starting a new course (step is ASK_TOPIC), discard any old slot values
        is_new_session = len(req.messages) <= 2 or all(m.get("content", "").lower().strip() in ["hi", "hello", "hey", "restart", "start", "create a new course"] for m in req.messages if m.get("role") == "user")
        if req.currentStep == "ASK_TOPIC" and is_new_session:
            current_slots = {
                "topic": "",
                "learningGoal": "",
                "currentLevel": "",
                "learningStyle": "",
                "duration": "",
                "language": "English"
            }

        # Check if the user is confirming/cancelling a topic change warning
        is_topic_confirm = False
        is_topic_cancel = False
        
        lowercase_msg = user_message.lower()
        confirm_words = ["yes", "continue", "looks good", "proceed", "generate", "correct", "confirm", "happy", "fine", "ok", "go ahead"]

        if req.messages:
            last_assistant_msg = ""
            for msg in reversed(req.messages):
                if msg.get("role") == "assistant":
                    last_assistant_msg = msg.get("content", "").lower()
                    break
            if "will require regenerating the syllabus outline" in last_assistant_msg:
                # User is replying to topic change warning
                if any(w in lowercase_msg for w in confirm_words) or "change" in lowercase_msg:
                    is_topic_confirm = True
                elif any(w in lowercase_msg for w in ["no", "cancel", "back", "stop"]):
                    is_topic_cancel = True

        clear_course_data_flag = False

        if is_topic_cancel:
            old_topic = details.get("topic") or details.get("subject") or details.get("courseName") or ""
            if old_topic:
                current_slots["topic"] = old_topic
            
            # Route back to previous state
            has_content = len(req.courseData.get("content", [])) > 0
            if has_content:
                meta = {
                    "next_step": "CONFIRM_GENERATE",
                    **current_slots
                }
                return JSONResponse({
                    "status": "success",
                    "reply": "Topic change cancelled. The course structure has been finalized. Would you like me to start generating the complete course content?",
                    "quickReplies": ["Yes, start generating", "No, go back to outline"],
                    "metadata": meta,
                    "type": "details"
                })
            else:
                meta = {
                    "next_step": "CONFIRM_DETAILS",
                    **current_slots
                }
                return JSONResponse({
                    "status": "success",
                    "reply": "Topic change cancelled. Here are your course details. Would you like to review the details again or proceed with the course structure?",
                    "quickReplies": ["Confirm details & proceed", "Edit Details"],
                    "metadata": meta,
                    "type": "details_card"
                })

        pending_topic = details.get("pending_topic")
        if is_topic_confirm and pending_topic:
            if pending_topic == "CLEAR":
                current_slots["topic"] = None
            else:
                current_slots["topic"] = pending_topic
            
            # Clear all other details slots to force re-answering them
            current_slots["learningGoal"] = ""
            current_slots["currentLevel"] = ""
            current_slots["learningStyle"] = ""
            current_slots["duration"] = ""
            
            # Also clear the existing structure and content to force regeneration!
            req.courseData["structure"] = {"modules": []}
            req.courseData["content"] = []
            clear_course_data_flag = True

        # Check if topic change request is triggered
        has_existing_structure = len(req.courseData.get("structure", {}).get("modules", [])) > 0
        is_requesting_topic_change = False
        
        # Temp NLU run to see if they specified a new topic in message
        temp_slots = {**current_slots}
        temp_updated, _ = extract_slots_from_message(user_message, temp_slots, req.currentStep, req.draft_id)
        old_topic_temp = current_slots.get("topic")
        new_topic_temp = temp_updated.get("topic")
        
        if old_topic_temp and new_topic_temp and str(old_topic_temp).lower().strip() != str(new_topic_temp).lower().strip():
            if has_existing_structure and not is_topic_confirm and not is_topic_cancel:
                is_requesting_topic_change = True
                
        if has_existing_structure and not is_topic_confirm and not is_topic_cancel:
            if any(w in lowercase_msg for w in ["change topic", "change the topic", "different topic", "another topic", "edit topic", "choose topic", "edit subject", "change subject"]):
                is_requesting_topic_change = True
                
        if is_requesting_topic_change:
            pending = new_topic_temp if (new_topic_temp and str(old_topic_temp).lower().strip() != str(new_topic_temp).lower().strip()) else "CLEAR"
            meta = {
                "next_step": "CONFIRM_DETAILS",
                "pending_topic": pending,
                **current_slots
            }
            if pending != "CLEAR":
                reply = f"Changing the course topic from '{old_topic_temp}' to '{pending}' will require regenerating the syllabus outline and course content from scratch. Are you sure you want to proceed with this change?"
            else:
                reply = f"Changing the course topic will require regenerating the syllabus outline and course content from scratch. Are you sure you want to proceed?"
            return JSONResponse({
                "status": "success",
                "reply": reply,
                "quickReplies": ["Yes, change topic", "Cancel"],
                "metadata": meta,
                "type": "details_card"
            })

        # Stage 1: Run NLU Slot Extraction
        updated_slots, raw_extracted = extract_slots_from_message(user_message, current_slots, req.currentStep, req.draft_id)
        
        if is_topic_confirm:
            updated_slots["pending_topic"] = None
        if is_topic_cancel:
            updated_slots["pending_topic"] = None

        # Check if topic changed and we already have structure/content
        old_topic = current_slots.get("topic")
        new_topic = updated_slots.get("topic")
        
        # If topic changed, clear all subsequent details slots
        if old_topic and new_topic and str(old_topic).lower().strip() != str(new_topic).lower().strip():
            if not has_existing_structure or is_topic_confirm:
                updated_slots["learningGoal"] = ""
                updated_slots["currentLevel"] = ""
                updated_slots["learningStyle"] = ""
                updated_slots["duration"] = ""
                req.courseData["structure"] = {"modules": []}
                req.courseData["content"] = []
                clear_course_data_flag = True

        # Stage 2: Programmatic Dialog Solver
        next_step, validation_error = determine_next_step(
            req.currentStep, 
            updated_slots, 
            user_message, 
            raw_extracted,
            has_existing_structure=has_existing_structure
        )

        # Check if the user is confirming/cancelling existing structure reuse
        is_structure_confirm = False
        is_structure_regenerate = False
        
        if req.messages:
            last_assistant_msg = ""
            for msg in reversed(req.messages):
                if msg.get("role") == "assistant":
                    last_assistant_msg = msg.get("content", "").lower()
                    break
            if "would you like to proceed with your existing outline structure" in last_assistant_msg:
                if any(w in lowercase_msg for w in ["keep", "existing", "proceed", "old", "yes"]):
                    is_structure_confirm = True
                elif any(w in lowercase_msg for w in ["generate", "new", "recreate", "change", "fresh"]):
                    is_structure_regenerate = True

        # Check if details changed (excluding topic, which is handled separately)
        details_changed_chk = False
        if details:
            for k in ["learningGoal", "currentLevel", "learningStyle", "duration"]:
                old_val = details.get(k) or ""
                new_val = updated_slots.get(k) or ""
                if str(old_val).lower().strip() != str(new_val).lower().strip():
                    details_changed_chk = True
                    break
                    
        has_existing_structure = len(req.courseData.get("structure", {}).get("modules", [])) > 0
        
        if next_step == "OUTLINE_EDIT" and details_changed_chk and has_existing_structure and not is_structure_confirm and not is_structure_regenerate:
            # Show selection prompt
            return JSONResponse({
                "status": "success",
                "reply": "You already have a generated module structure. Would you like to proceed with your existing outline structure or generate a new one based on the updated details?",
                "quickReplies": ["Keep existing structure", "Generate new structure"],
                "metadata": {
                    "next_step": "CONFIRM_DETAILS",
                    "pending_details_change": "true",
                    **updated_slots
                },
                "type": "details_card"
            })

        if is_structure_regenerate:
            req.courseData["structure"] = {"modules": []}
            req.courseData["content"] = []

        if next_step == "OUTLINE_EDIT" and req.currentStep in ["CONFIRM_DETAILS", "EDIT_DETAILS_CHOICE"]:
            has_content = len(req.courseData.get("content", [])) > 0
            if has_content and (is_structure_confirm or not details_changed_chk):
                next_step = "CONFIRM_GENERATE"

        # Stage 3: Dynamic NLG Prompt Generation
        system_prompt = build_builder_system_prompt(next_step, updated_slots, validation_error)

        cleaned_history = reinject_quick_replies_into_history(req.messages, updated_slots)
        messages = [{"role": "system", "content": system_prompt}] + cleaned_history

        response = openai_client.chat.completions.create(
            model=LLM_MODEL,
            messages=messages,
            temperature=0.7,
            max_tokens=4000
        )

        from metering_helper import track_chatbot_cost
        track_chatbot_cost(req.draft_id, response, LLM_MODEL, f"chatbot_chat_{next_step}")

        ai_reply = response.choices[0].message.content

        # Determine scope based on next_step
        scope = "Details"
        if next_step == "PROMPT_GEN":
            scope = "Content"
        elif next_step == "OUTLINE_EDIT":
            scope = "Structure"

        # Parse quick-replies lists first to prevent clean_reply_text in parse_metadata from stripping them
        ai_reply, quick_replies = parse_quick_replies(ai_reply)

        # If LLM failed to output quick replies, apply fallback safety net choices based on next_step
        if not quick_replies:
            topic_lower = str(updated_slots.get("topic", "")).lower()
            is_programming = any(x in topic_lower for x in ["python", "java", "c++", "coding", "program", "developer", "react", "javascript", "typescript", "sql", "backend", "frontend", "software", "git", "c#", "html", "css", "database", "node", "express"])
            
            if next_step == "ASK_TOPIC":
                quick_replies = ["Python Programming", "English Grammar", "Digital Marketing", "Machine Learning"]
            elif next_step == "ASK_GOAL":
                quick_replies = ["Build a Web App", "Automate Excel Tasks", "Data Analysis & AI", "Get a Developer Job"]
            elif next_step == "ASK_LEVEL":
                quick_replies = ["Complete Beginner / Start Fresh", "Intermediate / Some experience", "Advanced / Deep Dive"]
            elif next_step == "ASK_STYLE":
                if is_programming:
                    quick_replies = ["Hands-on Coding", "Interactive Quizzes", "Detailed Explanations", "Balanced Combination"]
                else:
                    quick_replies = ["Detailed Explanations", "Interactive Quizzes", "Structured Tables", "Balanced Combination"]
            elif next_step == "ASK_DURATION":
                quick_replies = ["1 Hour", "2 Hours", "5 Hours", "10 Hours", "15 Hours", "20 Hours"]
            elif next_step == "CONFIRM_DETAILS":
                quick_replies = ["Confirm details & proceed", "Change topic", "Change duration", "Change level"]
            elif next_step == "EDIT_DETAILS_CHOICE":
                quick_replies = ["Edit Topic", "Edit Learning Goal", "Edit Difficulty Level", "Edit Learning Style", "Edit Duration"]
            elif next_step == "ASK_GENERATE_SKELETON":
                quick_replies = ["Yes, generate modules!", "Go back"]
            elif next_step == "ASK_REDUCE_COUNT":
                quick_replies = ["Reduce by 1 module", "Reduce by 2 modules", "Your choice (Reduce by 2)"]
            elif next_step == "ASK_ADD_TOPIC":
                quick_replies = ["Your choice", "Add specific topic"]
            elif next_step == "OUTLINE_EDIT":
                quick_replies = ["Confirm Outline", "Reduce modules", "Add new module", "Rename modules/chapters"]
                current_modules = req.courseData.get("structure", {}).get("modules", [])
                if len(current_modules) <= 1:
                    quick_replies = ["Confirm Outline", "Add new module", "Rename modules/chapters"]
            elif next_step == "EDIT_OUTLINE_CHOICE":
                quick_replies = ["Reduce modules", "Add new module", "Rename modules/chapters", "Reorder modules"]
                current_modules = req.courseData.get("structure", {}).get("modules", [])
                if len(current_modules) <= 1:
                    quick_replies = ["Add new module", "Rename modules/chapters", "Reorder modules"]
            elif next_step == "CONFIRM_GENERATE":
                quick_replies = ["Yes, generate content", "No, go back to outline"]

        # Parse metadata suggestions
        reply_text, metadata, type_val = parse_metadata(
            ai_reply=ai_reply,
            scope=scope,
            details=updated_slots
        )

        # Overwrite or inject next_step and clean slots into details metadata block
        if next_step == "PROMPT_GEN":
            type_val = "content"
            if metadata:
                metadata["next_step"] = "CONFIRM_GENERATE"
            else:
                metadata = {
                    "next_step": "CONFIRM_GENERATE",
                    "prompts": []
                }
        elif next_step == "OUTLINE_EDIT":
            type_val = "structure"
            
            # If details changed (e.g. topic, goal, level, style, duration), discard the old structure to trigger regeneration
            details = req.courseData.get("details", {})
            details_changed = False
            if is_structure_regenerate:
                details_changed = True
            elif is_structure_confirm:
                details_changed = False
            else:
                if details:
                    old_topic = details.get("topic") or details.get("subject") or details.get("courseName") or ""
                    new_topic = updated_slots.get("topic") or updated_slots.get("subject") or updated_slots.get("courseName") or ""
                    if str(old_topic).lower().strip() != str(new_topic).lower().strip():
                        details_changed = True
                        
                    old_goal = details.get("learningGoal") or details.get("description") or details.get("goal") or ""
                    new_goal = updated_slots.get("learningGoal") or updated_slots.get("description") or updated_slots.get("goal") or ""
                    if str(old_goal).lower().strip() != str(new_goal).lower().strip():
                        details_changed = True
                        
                    old_level = details.get("currentLevel") or details.get("level") or ""
                    new_level = updated_slots.get("currentLevel") or updated_slots.get("level") or ""
                    if str(old_level).lower().strip() != str(new_level).lower().strip():
                        details_changed = True
                        
                    old_style = details.get("learningStyle") or details.get("requirements") or details.get("style") or ""
                    new_style = updated_slots.get("learningStyle") or updated_slots.get("requirements") or updated_slots.get("style") or ""
                    if str(old_style).lower().strip() != str(new_style).lower().strip():
                        details_changed = True
                        
                    old_duration = details.get("duration") or details.get("courseDuration") or details.get("hours") or ""
                    new_duration = updated_slots.get("duration") or updated_slots.get("courseDuration") or updated_slots.get("hours") or ""
                    if str(old_duration).lower().strip() != str(new_duration).lower().strip():
                        details_changed = True

            current_structure_modules = req.courseData.get("structure", {}).get("modules", [])
            if details_changed:
                logger.info("Course details changed by user. Discarding old outline structure for regeneration.")
                current_structure_modules = []
                
            if not current_structure_modules:
                try:
                    logger.info("Generating initial syllabus structure using generate_course_structure...")
                    structure_res = generate_course_structure(
                        courseName=updated_slots.get("topic") or "Custom Course",
                        description=updated_slots.get("learningGoal") or "Course Content",
                        subject=updated_slots.get("topic") or "General",
                        level=updated_slots.get("currentLevel") or "beginner"
                    )
                    modules = structure_res.get("modules", [])
                    metadata = {
                        "next_step": "OUTLINE_EDIT",
                        "modules": modules
                    }
                except Exception as structure_err:
                    logger.error(f"Failed to generate initial structure on the fly: {structure_err}")
                    metadata = {
                        "next_step": "OUTLINE_EDIT",
                        "modules": []
                    }
            else:
                if not metadata:
                    metadata = {}
                metadata["next_step"] = "OUTLINE_EDIT"
                metadata["modules"] = current_structure_modules
        else:
            if next_step in ["CONFIRM_DETAILS", "EDIT_DETAILS_CHOICE"]:
                type_val = "details_card"
            else:
                type_val = "details"
            if metadata:
                metadata["next_step"] = next_step
                for k, v in updated_slots.items():
                    metadata[k] = v
            else:
                # Minimal fallback metadata so frontend receives step transitions & slot values
                metadata = {
                    "next_step": next_step,
                    **updated_slots
                }

        # Quick replies are already parsed and extracted above
        reply_text = reply_text.strip()

        logger.info(f"[Chatbot Solver] Step transition: {req.currentStep} -> {next_step} | Metadata: {metadata}")

        if clear_course_data_flag:
            if not metadata:
                metadata = {}
            metadata["clear_course_data"] = True

        return {
            "status": "success",
            "reply": reply_text,
            "quickReplies": quick_replies,
            "metadata": metadata,
            "type": type_val
        }
    except Exception as e:
        logger.error(f"Error in chatbot builder completion API: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


from schemas import ChatbotDraftSaveRequest

# Draft database table startup init event
@app.on_event("startup")
def startup_db_init():
    try:
        from database import init_draft_table
        init_draft_table()
    except Exception as e:
        logger.error(f"Failed to auto-initialize drafts MySQL table: {e}")


# Active background generation tasks registry
bg_generation_registry = {}

from pydantic import BaseModel
import threading

class StartBgGenRequest(BaseModel):
    draft_id: str
    courseData: dict
    messages: list

def run_background_generation(draft_id: str, course_data: dict, messages: list):
    try:
        from database import save_chatbot_draft
        from app.api.online_course_generator import generate_lesson_blocks
        from schemas import LessonRequest, CourseDetails
        import asyncio
        import json
        import re

        modules = course_data.get("structure", {}).get("modules", [])
        chapters_to_generate = []
        for m_idx, m in enumerate(modules):
            module_title = m.get("title", "")
            chapters = m.get("chapters", [])
            for c_idx, c in enumerate(chapters):
                chapter_title = c.get("title", "")
                contents = c.get("contents", [])
                has_content = any(item.get("type") == "lesson-blocks" for item in contents)
                if not has_content:
                    chapters_to_generate.append({
                        "m_idx": m_idx,
                        "c_idx": c_idx,
                        "module_title": module_title,
                        "chapter_title": chapter_title,
                        "prompt": c.get("prompt") or f"Generate a detailed structured lesson on {chapter_title}"
                    })

        total = len(chapters_to_generate)
        if total == 0:
            bg_generation_registry[draft_id] = {
                "status": "completed",
                "completed": 0,
                "total": 0,
                "current_title": ""
            }
            return

        bg_generation_registry[draft_id] = {
            "status": "generating",
            "completed": 0,
            "total": total,
            "current_title": chapters_to_generate[0]["chapter_title"],
            "cancel_requested": False
        }

        details = course_data.get("details", {})
        course_details_obj = CourseDetails(
            courseType=details.get("courseType") or "Custom Course",
            subject=details.get("subject") or "",
            courseName=details.get("courseName") or "",
            description=details.get("description") or "",
            price=details.get("price") or "",
            duration=details.get("duration") or "",
            requirements=details.get("requirements") or "",
            level=details.get("level") or "beginner",
            language=details.get("language") or "English",
            scriptingLanguage=details.get("scriptingLanguage") or "NA",
            evaluator=details.get("evaluator") or ""
        )

        updated_course = json.loads(json.dumps(course_data))

        for i, ch in enumerate(chapters_to_generate):
            if bg_generation_registry.get(draft_id, {}).get("cancel_requested"):
                bg_generation_registry[draft_id]["status"] = "cancelled"
                logger.info(f"[BG Generation] Task {draft_id} cancelled.")
                return

            m_idx = ch["m_idx"]
            c_idx = ch["c_idx"]
            chapter_title = ch["chapter_title"]
            module_title = ch["module_title"]
            prompt = ch["prompt"]

            bg_generation_registry[draft_id]["current_title"] = chapter_title
            bg_generation_registry[draft_id]["completed"] = i

            req_obj = LessonRequest(
                title=chapter_title,
                module_title=module_title,
                prompt=prompt,
                type="html",
                course_details=course_details_obj,
                draft_id=draft_id
            )

            # Run async function in a new loop
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                res_block = loop.run_until_complete(generate_lesson_blocks(req_obj))
            finally:
                loop.close()

            if res_block and res_block.blocks:
                res_block_dict = res_block.dict() if hasattr(res_block, "dict") else res_block
                blocks_data = res_block_dict.get("blocks", [])
                latest_modules = updated_course["structure"]["modules"]
                target_chapter = latest_modules[m_idx]["chapters"][c_idx]
                target_chapter["contents"] = [{
                    "type": "lesson-blocks",
                    "title": res_block_dict.get("title") or chapter_title,
                    "blocks": blocks_data,
                    "source": "ai",
                    "completed": True,
                    "timestamp": ""
                }]
                target_chapter["content"] = {
                    "content_type": "lesson-blocks",
                    "html_content": "",
                    "completed": True
                }

                if "content" not in updated_course:
                    updated_course["content"] = []
                
                content_exist = False
                for idx, item in enumerate(updated_course["content"]):
                    if item.get("module_title") == module_title and item.get("chapter_title") == chapter_title:
                        updated_course["content"][idx] = {
                            "module_title": module_title,
                            "chapter_title": chapter_title,
                            "prompt": prompt,
                            "type": "html",
                            "blocks": blocks_data
                        }
                        content_exist = True
                        break
                if not content_exist:
                    updated_course["content"].append({
                        "module_title": module_title,
                        "chapter_title": chapter_title,
                        "prompt": prompt,
                        "type": "html",
                        "blocks": blocks_data
                    })

            # Save draft database
            save_chatbot_draft(
                draft_id=draft_id,
                course_name=updated_course.get("details", {}).get("courseName") or "Custom Course",
                current_step="CONFIRM_GENERATE",
                course_data=updated_course,
                messages=messages
            )
            bg_generation_registry[draft_id]["completed"] = i + 1

        bg_generation_registry[draft_id]["status"] = "completed"
        # Final transition to READY
        save_chatbot_draft(
            draft_id=draft_id,
            course_name=updated_course.get("details", {}).get("courseName") or "Custom Course",
            current_step="READY",
            course_data=updated_course,
            messages=messages
        )
        logger.info(f"[BG Generation] Task {draft_id} completed successfully.")

    except Exception as ex:
        logger.error(f"[BG Generation] Task {draft_id} failed: {ex}")
        if draft_id in bg_generation_registry:
            bg_generation_registry[draft_id]["status"] = "failed"

@app.post("/course/chatbot-builder/generate-content/start")
async def api_start_content_generation(req: StartBgGenRequest):
    draft_id = req.draft_id
    if draft_id in bg_generation_registry and bg_generation_registry[draft_id]["status"] == "generating":
        return {"status": "success", "message": "Generation already running"}
    t = threading.Thread(
        target=run_background_generation,
        args=(draft_id, req.courseData, req.messages),
        daemon=True
    )
    t.start()
    return {"status": "success", "message": "Background generation started"}

@app.get("/course/chatbot-builder/generate-content/status/{draft_id}")
async def api_get_generation_status(draft_id: str):
    if draft_id not in bg_generation_registry:
        from database import get_chatbot_draft
        try:
            draft = get_chatbot_draft(draft_id)
            if draft and (draft.get("current_step") == "READY" or draft.get("currentStep") == "READY"):
                return {
                    "status": "completed",
                    "completed": 100,
                    "total": 100,
                    "current_title": ""
                }
        except Exception:
            pass
        return {"status": "idle", "completed": 0, "total": 0, "current_title": ""}
    status_info = bg_generation_registry[draft_id]
    return {
        "status": status_info["status"],
        "completed": status_info["completed"],
        "total": status_info["total"],
        "current_title": status_info["current_title"]
    }

@app.post("/course/chatbot-builder/generate-content/cancel/{draft_id}")
async def api_cancel_generation(draft_id: str):
    if draft_id in bg_generation_registry:
        bg_generation_registry[draft_id]["cancel_requested"] = True
        bg_generation_registry[draft_id]["status"] = "cancelled"
        return {"status": "success", "message": "Cancellation requested"}
    return {"status": "error", "message": "No active generation task found for this draft"}

# Get all drafts
@app.get("/course/chatbot-builder/drafts")
def api_get_chatbot_drafts():
    try:
        from database import get_chatbot_drafts
        drafts = get_chatbot_drafts()
        return {"status": "success", "drafts": drafts}
    except Exception as e:
        logger.error(f"Error fetching chatbot drafts: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Get a specific draft
@app.get("/course/chatbot-builder/draft/{draft_id}")
def api_get_chatbot_draft(draft_id: str):
    try:
        from database import get_chatbot_draft
        draft = get_chatbot_draft(draft_id)
        if not draft:
            raise HTTPException(status_code=404, detail="Draft not found")
        return {"status": "success", "draft": draft}
    except HTTPException as he:
        raise he
    except Exception as e:
        logger.error(f"Error fetching chatbot draft {draft_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Save or update a draft
@app.post("/course/chatbot-builder/draft")
def api_save_chatbot_draft(req: ChatbotDraftSaveRequest):
    try:
        from database import save_chatbot_draft
        save_chatbot_draft(
            draft_id=req.id,
            course_name=req.courseName,
            current_step=req.currentStep,
            course_data=req.courseData,
            messages=req.messages
        )
        return {"status": "success", "message": "Draft saved successfully"}
    except Exception as e:
        logger.error(f"Error saving chatbot draft: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# Delete a draft
@app.delete("/course/chatbot-builder/draft/{draft_id}")
def api_delete_chatbot_draft(draft_id: str):
    try:
        from database import delete_chatbot_draft
        delete_chatbot_draft(draft_id)
        return {"status": "success", "message": "Draft deleted successfully"}
    except Exception as e:
        logger.error(f"Error deleting chatbot draft {draft_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

from schemas import RenameDraftRequest

# Rename a draft
@app.post("/course/chatbot-builder/draft/{draft_id}/rename")
def api_rename_chatbot_draft(draft_id: str, req: RenameDraftRequest):
    try:
        from database import rename_chatbot_draft
        rename_chatbot_draft(draft_id, req.name)
        return {"status": "success", "message": "Draft renamed successfully"}
    except Exception as e:
        logger.error(f"Error renaming chatbot draft {draft_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))






