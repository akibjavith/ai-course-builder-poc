import fitz  # PyMuPDF
import docx
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from urllib.parse import urlparse, parse_qs

def extract_text_from_pdf(filepath: str) -> str:
    doc = fitz.open(filepath)
    text = ""
    for page in doc:
        text += page.get_text() + "\n"
    return text.strip()

def extract_text_from_docx(filepath: str) -> str:
    doc = docx.Document(filepath)
    return "\n".join([p.text for p in doc.paragraphs]).strip()

def extract_text_from_pptx(filepath: str) -> str:
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_text_from_url(url: str) -> str:
    headers = {"User-Agent": "Mozilla/5.0"}
    res = requests.get(url, headers=headers, timeout=10)
    res.raise_for_status()
    soup = BeautifulSoup(res.text, "html.parser")
    # Kill script and style elements
    for script in soup(["script", "style"]):
        script.decompose()
    return soup.get_text(separator="\n").strip()

def extract_youtube_transcript(url: str) -> str:
    parsed_url = urlparse(url)
    video_id = ""
    if "youtube.com" in parsed_url.netloc:
        video_id = parse_qs(parsed_url.query).get("v", [""])[0]
    elif "youtu.be" in parsed_url.netloc:
        video_id = parsed_url.path.lstrip("/")
        
    if not video_id:
        raise ValueError("Invalid YouTube URL")
        
    transcript_list = YouTubeTranscriptApi.get_transcript(video_id)
    text = " ".join([t['text'] for t in transcript_list])
    return text.strip()
